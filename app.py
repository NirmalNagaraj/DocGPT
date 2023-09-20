from dotenv import load_dotenv
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from langchain.callbacks import get_openai_callback
from docx import Document
from pptx import Presentation
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns

def extract_column_name(prompt):
    keywords = ["of", "on", "from"]
    words = prompt.split()
    for i, word in enumerate(words):
        if word in keywords and i < len(words) - 1:
            return words[i + 1]
    return None

def extract_person_name(query):
    words = query.split()
    for i, word in enumerate(words):
        if word.lower() == "name" and i < len(words) - 1:
            return words[i + 1]
    return None

def main():
    load_dotenv()
    st.set_page_config(
        page_title="DocGPT",
        page_icon="📚",
        initial_sidebar_state="expanded",
    )
   
    st.markdown(
        """
        <style>
        .stApp header {
            background-color: #000000;
        }
        .stApp header a, .stApp header h1, .stApp header h2 {
            color: black;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )
    st.header("DocGPT 📚")

    uploaded_file = st.file_uploader("Upload your document", type=["pdf", "docx", "txt", "pptx", "csv", "xlsx"], key="file_uploader")

    if uploaded_file is not None:
        if uploaded_file.type == 'application/pdf':
            pdf_reader = PdfReader(uploaded_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
        elif uploaded_file.type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            doc = Document(uploaded_file)
            text = "\n".join([para.text for para in doc.paragraphs])
        elif uploaded_file.type == 'text/plain':
            text = uploaded_file.read().decode('utf-8')
        elif uploaded_file.type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            ppt = Presentation(uploaded_file)
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
        elif uploaded_file.type == 'text/csv':
            df = pd.read_csv(uploaded_file)
            st.write("Uploaded CSV Data:")
            st.write(df)

            prompt = st.text_input("Enter a data science analysis prompt:")

            if prompt:
                st.write("Analysis Result:")
                prompt_lower = prompt.lower()
                column = extract_column_name(prompt)
                if column is not None:
                    if "mean" in prompt_lower or "average" in prompt_lower:
                        st.write(f"Mean of '{column}': {df[column].mean()}")

                    elif "median" in prompt_lower:
                        st.write(f"Median of '{column}': {df[column].median()}")

                    elif "mode" in prompt_lower:
                        mode_values = df[column].mode()
                        if not mode_values.empty:
                            st.write(f"Mode of '{column}': {', '.join(map(str, mode_values))}")
                        else:
                            st.write(f"No mode found in '{column}'.")

                    elif "histogram" in prompt_lower:
                        plt.hist(df[column], bins=20)
                        st.pyplot(plt)

                    elif "scatterplot" in prompt_lower:
                        x_column = st.selectbox("Select the X-axis column:", df.columns)
                        plt.scatter(df[x_column], df[column])
                        st.pyplot(plt)

                    elif "count" in prompt_lower:
                        st.write(f"Count of '{column}': {df[column].count()}")

                    elif "sum" in prompt_lower:
                        st.write(f"Sum of '{column}': {df[column].sum()}")

                    elif "null" in prompt_lower:
                        st.write(f"Null value count in '{column}': {df[column].isnull().sum()}")

                    elif "min" in prompt_lower:
                        st.write(f"Min value in '{column}': {df[column].min()}")

                    elif "max" in prompt_lower:
                        st.write(f"Max value in '{column}': {df[column].max()}")

                    elif "line plot" in prompt_lower:
                        x_column = st.selectbox("Select the X-axis column:", df.columns)
                        y_column = st.selectbox("Select the Y-axis column:", df.columns)
                        plt.plot(df[x_column], df[y_column])
                        st.pyplot(plt)

                    elif "scatter chart" in prompt_lower:
                        x_column = st.selectbox("Select the X-axis column:", df.columns)
                        y_column = st.selectbox("Select the Y-axis column:", df.columns)
                        plt.scatter(df[x_column], df[y_column])
                        st.pyplot(plt)

                    elif "correlation chart" in prompt_lower:
                        corr_matrix = df.corr()
                        sns.heatmap(corr_matrix, annot=True)
                        st.pyplot(plt)

                    elif "heatmap" in prompt_lower:
                        sns.heatmap(df.corr(), annot=True, cmap="coolwarm")
                        st.pyplot(plt)

                    elif "bubble chart" in prompt_lower:
                        x_column = st.selectbox("Select the X-axis column:", df.columns)
                        y_column = st.selectbox("Select the Y-axis column:", df.columns)
                        size_column = st.selectbox("Select the size column:", df.columns)
                        plt.scatter(df[x_column], df[y_column], s=df[size_column])
                        st.pyplot(plt)

                    elif "radar chart" in prompt_lower:
                        # Implement radar chart logic here
                        pass

                    elif "ridge plot" in prompt_lower:
                        sns.ridgeplot(df)
                        st.pyplot(plt)

                    elif "dendrogram" in prompt_lower:
                        corr_matrix = df.corr()
                        linkage_matrix = hierarchy.linkage(corr_matrix, method='ward')
                        dendrogram = hierarchy.dendrogram(linkage_matrix, labels=corr_matrix.index)
                        st.pyplot(plt)

                    else:
                        st.write("Unsupported analysis prompt.")

                else:
                    st.write("Column not specified in the prompt.")

            return
        elif uploaded_file.type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            df = pd.read_excel(uploaded_file, engine="openpyxl")
            st.write("Uploaded Excel Data:")
            st.write(df)

            user_prompt = st.text_input("Enter a prompt for Excel analysis:")

            if user_prompt:
                st.write("Excel Analysis Result:")
                user_prompt_lower = user_prompt.lower()
                column = extract_column_name(user_prompt)
                if column is not None:
                    if "mean" in user_prompt_lower or "average" in user_prompt_lower:
                        st.write(f"Mean of '{column}': {df[column].mean()}")

                    # Add more analysis options here based on user prompts

                    else:
                        st.write("Unsupported analysis prompt.")

                else:
                    st.write("Column not specified in the prompt.")

            return
        else:
            st.warning("Unsupported file type. Please upload a PDF, DOCX, TXT, PPTX, XLSX, or CSV file.")
            return

        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=2000,
            chunk_overlap=400,
            length_function=len
        )
        chunks = text_splitter.split_text(text)

        embeddings = OpenAIEmbeddings()
        knowledge_base = FAISS.from_texts(chunks, embedding=embeddings)

        user_question = st.text_input("Ask a question about your document:")
        if user_question:
            docs = knowledge_base.similarity_search(user_question)

            llm = OpenAI()
            chain = load_qa_chain(llm, chain_type="stuff")
            with get_openai_callback() as cb:
                response = chain.run(input_documents=docs, question=user_question)

            st.write(response)

if __name__ == '__main__':
    main()
